"""Module for converting presentation files to Markdown format."""

import logging
from pathlib import Path
from typing import List
try:
    from pptx import Presentation
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

from .base import BaseConverter

logger = logging.getLogger(__name__)


class PresentationConverter(BaseConverter):
    """Converter for presentation files to Markdown format."""
    
    def convert(self, file_path: Path, output_dir: Path) -> List[Path]:
        """
        Convert a presentation file to Markdown format.
        
        Args:
            file_path (Path): Path to the presentation file
            output_dir (Path): Directory to save the output files
            
        Returns:
            List[Path]: List of paths to the converted Markdown files
        """
        if not HAS_PYTHON_PPTX:
            raise ImportError("python-pptx is required for presentation conversion. Please install it with 'pip install python-pptx'")
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Create output directory if it doesn't exist
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Load the presentation
        presentation = Presentation(file_path)
        
        # Convert to Markdown
        md_content = self._presentation_to_markdown(presentation)
        
        # Save to file
        output_file = output_dir / f"{file_path.stem}.md"
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(md_content)
            
        logger.info(f"Converted {file_path.name} to {output_file}")
        return [output_file]
    
    def _presentation_to_markdown(self, presentation) -> str:
        """Convert presentation to Markdown format."""
        md_lines = []
        
        # Add title
        md_lines.append(f"# Presentation: {presentation.core_properties.title or 'Untitled'}")
        md_lines.append("")
        
        # Add author if available
        if presentation.core_properties.author:
            md_lines.append(f"**Author:** {presentation.core_properties.author}")
            md_lines.append("")
        
        # Process each slide
        for i, slide in enumerate(presentation.slides, 1):
            md_lines.append(f"## Slide {i}")
            md_lines.append("")
            
            # Process text content
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Try to determine if it's a title
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        md_lines.append(f"### {shape.text.strip()}")
                    else:
                        md_lines.append(shape.text.strip())
                    md_lines.append("")
            
            md_lines.append("---")  # Separator between slides
            md_lines.append("")
            
        return "\n".join(md_lines)
    
    def supported_extensions(self) -> List[str]:
        """
        Get list of supported file extensions.
        
        Returns:
            List[str]: List of supported extensions
        """
        return ['.pptx', '.ppt']